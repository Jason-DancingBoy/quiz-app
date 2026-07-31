import os
import re
from dataclasses import dataclass

from backend.logger import get_logger

logger = get_logger(__name__)


@dataclass
class ParsedDocument:
    title: str
    chunks: list[str]


def parse_text(text: str, title: str = "pasted") -> ParsedDocument:
    chunks = chunk_text(text)
    return ParsedDocument(title=title, chunks=chunks)


def parse_markdown(filepath: str) -> ParsedDocument:
    with open(filepath, "r", encoding="utf-8") as f:
        text = f.read()
    title = os.path.basename(filepath)
    return ParsedDocument(title=title, chunks=chunk_text(text))


def parse_pdf(filepath: str) -> ParsedDocument:
    text = ""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    except Exception:
        pass

    if not text.strip():
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(filepath)

    title = os.path.basename(filepath)
    return ParsedDocument(title=title, chunks=chunk_text(text))


def parse_docx(filepath: str) -> ParsedDocument:
    from docx import Document as DocxDocument
    doc = DocxDocument(filepath)
    text = "\n".join(p.text for p in doc.paragraphs)
    title = os.path.basename(filepath)
    return ParsedDocument(title=title, chunks=chunk_text(text))


def _split_long_paragraph(para: str, max_len: int) -> list[str]:
    """Split a long paragraph at sentence boundaries, falling back to hard cut."""
    sentences = re.split(r"(?<=[。！？；\n])", para)
    result = []
    current = ""
    for sent in sentences:
        if not sent.strip():
            continue
        if len(current) + len(sent) <= max_len:
            current += sent
        else:
            if current:
                result.append(current.strip())
            if len(sent) >= max_len:
                for i in range(0, len(sent), max_len):
                    result.append(sent[i:i + max_len].strip())
                current = ""
            else:
                current = sent
    if current.strip():
        result.append(current.strip())
    return result


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    """Split text by paragraph boundaries, falling back to sentence boundaries, then hard cut."""
    paragraphs = re.split(r"\n\s*\n", text.strip())
    chunks = []
    current = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current) + len(para) <= chunk_size:
            current += ("\n\n" if current else "") + para
        else:
            if current:
                chunks.append(current)
            if len(para) >= chunk_size:
                effective_size = chunk_size - overlap
                for sub_chunk in _split_long_paragraph(para, effective_size):
                    chunks.append(sub_chunk)
                current = ""
            else:
                current = para

    if current:
        chunks.append(current)

    if overlap > 0 and len(chunks) > 1:
        overlapped = [chunks[0]]
        for i in range(1, len(chunks)):
            prev_end = chunks[i - 1][-overlap:]
            overlapped.append(prev_end + "\n\n" + chunks[i])
        chunks = overlapped

    return chunks


def parse_pptx(filepath: str) -> ParsedDocument:
    from pptx import Presentation
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    title = os.path.basename(filepath)
    return ParsedDocument(title=title, chunks=chunk_text("\n".join(texts)))


PARSERS = {
    ".md": parse_markdown,
    ".txt": parse_markdown,
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}


def parse_file(filepath: str) -> ParsedDocument:
    ext = os.path.splitext(filepath)[1].lower()
    logger.info("Parsing file: %s (type=%s)", filepath, ext)
    parser = PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext}")
    result = parser(filepath)
    logger.info("File parsed successfully: %s (%d chunks)", filepath, len(result.chunks))
    return result
